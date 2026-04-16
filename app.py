from flask import Flask, request, send_file, jsonify, render_template
from pdf2docx import Converter
import subprocess, os, uuid, zipfile, time
from pathlib import Path

app = Flask(__name__)

UPLOAD = "uploads"
OUTPUT = "outputs"
PASSWORD = "kensony123"   # 🔐 CHANGE THIS

app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB limit

os.makedirs(UPLOAD, exist_ok=True)
os.makedirs(OUTPUT, exist_ok=True)

def uid(ext):
    return str(uuid.uuid4()) + ext

def check_auth():
    return request.headers.get("x-password") == PASSWORD

def cleanup():
    now = time.time()
    for folder in [UPLOAD, OUTPUT]:
        for f in Path(folder).glob("*"):
            if f.is_file() and now - f.stat().st_mtime > 3600:
                try:
                    f.unlink()
                except:
                    pass

# ── PDF → DOCX ──
@app.route("/convert/pdf-to-docx", methods=["POST"])
def pdf_to_docx():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    files = request.files.getlist("files")
    results = []

    for f in files:
        inp = os.path.join(UPLOAD, uid(".pdf"))
        out = os.path.join(OUTPUT, uid(".docx"))
        f.save(inp)

        cv = Converter(inp)
        cv.convert(out)
        cv.close()

        results.append(out)

    return _deliver(results, "converted_docx")


# ── DOCX → PDF ──
@app.route("/convert/docx-to-pdf", methods=["POST"])
def docx_to_pdf():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    files = request.files.getlist("files")
    results = []

    for f in files:
        inp = os.path.join(UPLOAD, uid(".docx"))
        f.save(inp)

        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", OUTPUT, inp
        ], check=True)

        results.append(os.path.join(OUTPUT, Path(inp).stem + ".pdf"))

    return _deliver(results, "converted_pdf")


# ── PPT → PDF ──
@app.route("/convert/ppt-to-pdf", methods=["POST"])
def ppt_to_pdf():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    files = request.files.getlist("files")
    results = []

    for f in files:
        ext = ".pptx" if f.filename.lower().endswith("pptx") else ".ppt"
        inp = os.path.join(UPLOAD, uid(ext))
        f.save(inp)

        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", OUTPUT, inp
        ], check=True)

        results.append(os.path.join(OUTPUT, Path(inp).stem + ".pdf"))

    return _deliver(results, "converted_pdf")


# ── PDF → PPT ──
@app.route("/convert/pdf-to-ppt", methods=["POST"])
def pdf_to_ppt():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    files = request.files.getlist("files")
    results = []

    for f in files:
        inp = os.path.join(UPLOAD, uid(".pdf"))
        out = os.path.join(OUTPUT, uid(".pptx"))
        f.save(inp)

        doc = fitz.open(inp)
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(UPLOAD, uid(".png"))
            pix.save(img_path)

            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(img_path, 0, 0,
                                     prs.slide_width, prs.slide_height)

        prs.save(out)
        results.append(out)

    return _deliver(results, "converted_pptx")


# ── MERGE PDFs ──
@app.route("/convert/merge-pdf", methods=["POST"])
def merge_pdf():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    import fitz

    files = request.files.getlist("files")
    if len(files) < 2:
        return jsonify({"error": "Upload at least 2 PDFs."}), 400

    merged = fitz.open()

    for f in files:
        inp = os.path.join(UPLOAD, uid(".pdf"))
        f.save(inp)
        merged.insert_pdf(fitz.open(inp))

    out = os.path.join(OUTPUT, uid(".pdf"))
    merged.save(out)

    return send_file(out, as_attachment=True, download_name="merged.pdf")


# ── AUDIO ──
@app.route("/convert/audio", methods=["POST"])
def audio_download():
    if not check_auth():
        return jsonify({"error": "Unauthorized"}), 403

    cleanup()
    data = request.get_json()
    url = (data or {}).get("url", "").strip()

    if not url:
        return jsonify({"error": "No URL provided."}), 400

    out_tmpl = os.path.join(OUTPUT, "%(title)s.%(ext)s")

    try:
        subprocess.run([
            "yt-dlp", "-x", "--audio-format", "mp3",
            "--audio-quality", "0",
            "-o", out_tmpl, url
        ], check=True, timeout=180)

        mp3s = sorted(Path(OUTPUT).glob("*.mp3"),
                      key=os.path.getmtime, reverse=True)

        if not mp3s:
            return jsonify({"error": "Download failed"}), 500

        return send_file(str(mp3s[0]), as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── ROOT ──
@app.route("/")
def index():
    return render_template("index.html")


# ── HELPER ──
def _deliver(paths, zip_name):
    valid = [p for p in paths if os.path.exists(p)]

    if not valid:
        return jsonify({"error": "No output."}), 500

    if len(valid) == 1:
        return send_file(valid[0], as_attachment=True)

    zip_path = os.path.join(OUTPUT, uid(".zip"))

    with zipfile.ZipFile(zip_path, "w") as zf:
        for p in valid:
            zf.write(p, Path(p).name)

    return send_file(zip_path, as_attachment=True,
                     download_name=zip_name + ".zip")


if __name__ == "__main__":
    app.run(debug=True, host="0.0.0.0", port=5000)